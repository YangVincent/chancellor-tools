#!/usr/bin/env python3
'''
This program pulls all presentations in a directory and merges them, adding the Delta Sigma Pi logo slide
in between each presentation.

The second version will have presentations sorted by order

The third version will have an agenda outputted in a .txt file

'''
from pptx import Presentation
from pptx.util import Inches
import os


def merge(pres): #open p1 and p2; merge
    prs = Presentation()
    img_path = 'logo.png'
    for each in pres:
        p = Presentation(pres)
        for slide in p.slides:
            prs.slides.add_slide(slide) #add each slide to main presentation
        
        #add divider slide
        slide = prs.slides.add_slide(prs.slide_layouts[6]) #blank slide
        left = top = Inches(1)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
    
   
    prs.save('gm.pptx')
    return prs
    

if __name__ == "__main__":
    #pull from directory
    files = [f for f in os.listdir('.') if os.path.isfile(f)]
    pres = []
    for f in files:
        if f[-5:] == ".pptx": #if file is a pptx 
        pres.append(f)
    
    #merge
    merge(pres)
